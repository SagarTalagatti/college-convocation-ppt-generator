from pptx import Presentation


def analyze_presentation():
    try:
        prs = Presentation("template.pptx")
        print("--- Analyzing 'template.pptx' ---")
        print(f"Found {len(prs.slide_layouts)} slide layouts.\n")

        for i, layout in enumerate(prs.slide_layouts):
            print("-----------------------------------------------------")
            print(f"Layout Index: {i}, Layout Name: '{layout.name}'")
            print("-----------------------------------------------------")

            # Add a temporary slide to inspect its runtime placeholders
            try:
                slide = prs.slides.add_slide(layout)

                if not slide.placeholders:
                    print("  This layout has no placeholders.")
                else:
                    print("  Placeholders found on this layout:")
                    for ph in slide.placeholders:
                        # We can get the type name directly from the enum
                        ph_type_name = "Unknown"
                        if ph.placeholder_format and ph.placeholder_format.type:
                            ph_type_name = ph.placeholder_format.type.name

                        print(
                            f"    -> Placeholder Index (idx): {ph.placeholder_format.idx}"
                        )
                        print(f"       Name: '{ph.name}'")
                        print(f"       Type: {ph_type_name}\n")

                # Clean up: remove the temporary slide
                element = slide.element
                parent = element.getparent()

                # --- THIS IS THE FIX ---
                # Only try to remove the element if its parent exists.
                # Some special layouts might not have a parent when
                # added this way, causing the 'NoneType' error.
                if parent is not None:
                    parent.remove(element)
                # --- END OF FIX ---

            except Exception as e:
                print(
                    f"    Could not add or inspect test slide for this layout. Error: {e}"
                )

            print("")

        print("--- Analysis Complete ---")
        print(
            "ACTION: Find your layout name (e.g., 'StudentLayout') in the list above."
        )
        print(
            "Note down the 'Placeholder Index (idx)' for your Name, USN, and Photo placeholders."
        )
        print("Update the 'create_presentation.py' script with these values.")

    except FileNotFoundError:
        print("Error: 'template.pptx' not found.")
        print(
            "Please create the template file first by following 'template_instructions.md'."
        )
    except Exception as e:
        print(f"An error occurred: {e}")


if __name__ == "__main__":
    analyze_presentation()
