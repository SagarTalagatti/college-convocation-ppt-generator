import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
import os

# --- 1. CONFIGURATION ---
# UPDATE THESE VALUES BASED ON YOUR 'find_layouts.py' OUTPUT

# The template file you created in PowerPoint
TEMPLATE_FILE = "template.pptx"

# The Excel file with student data
EXCEL_FILE = "sponsored_scholarships.xlsx"

# The name of the final presentation file that will be created
OUTPUT_FILE = "sponsored_scholarships_slides.pptx"

# The name of the custom slide layout you created (e.g., 'StudentLayout')
LAYOUT_NAME = "StudentLayout"

# The index (idx) of the placeholder for the STUDENT'S NAME
PLACEHOLDER_IDX_NAME = 11  # Example: 10

# The index (idx) of the placeholder for the STUDENT'S USN
PLACEHOLDER_IDX_USN = 12  # Example: 11

# The index (idx) of the placeholder for the STUDENT'S PHOTO
PLACEHOLDER_IDX_PHOTO = 10  # Example: 12

# --- 2. SCRIPT LOGIC (No edits needed below this line) ---


def load_data():
    """Loads student data from the specified Excel file."""
    try:
        df = pd.read_excel(EXCEL_FILE)
        # Ensure column names are what we expect
        expected_cols = ["name", "usn", "photo_path"]
        if not all(col in df.columns for col in expected_cols):
            print(
                f"Error: Excel file must contain columns: 'name', 'usn', 'photo_path'"
            )
            return None
        return df
    except FileNotFoundError:
        print(f"Error: Excel file not found at '{EXCEL_FILE}'")
        return None
    except Exception as e:
        print(f"Error reading Excel file: {e}")
        return None


def find_slide_layout(prs):
    """Finds the specified slide layout in the presentation."""
    for layout in prs.slide_layouts:
        if layout.name == LAYOUT_NAME:
            return layout
    print(f"Error: Slide layout '{LAYOUT_NAME}' not found in template.")
    print("Please check the LAYOUT_NAME variable or your template file.")
    print("Available layouts:")
    for layout in prs.slide_layouts:
        print(f"  - {layout.name}")
    return None


def create_presentation():
    """Main function to create the presentation."""

    # Check for placeholder configuration
    if "YOUR_LAYOUT_NAME_HERE" in LAYOUT_NAME:
        print("Error: Please update the configuration variables (LAYOUT_NAME, etc.)")
        print("in this script with the values from 'find_layouts.py' before running.")
        return

    print(f"Loading data from '{EXCEL_FILE}'...")
    df = load_data()
    if df is None:
        return

    try:
        prs = Presentation(TEMPLATE_FILE)
    except Exception as e:
        print(f"Error: Could not open template file '{TEMPLATE_FILE}'.")
        print(f"Details: {e}")
        return

    slide_layout = find_slide_layout(prs)
    if slide_layout is None:
        return

    print(f"Found layout '{LAYOUT_NAME}'. Starting to add slides...")

    # Loop through each row in the Excel file
    for index, row in df.iterrows():
        name = row["name"].upper()
        usn = row["usn"].upper()
        photo_path = row["photo_path"]

        print(f"  -> Processing: {name} ({usn})")

        # Add a new slide using our custom layout
        slide = prs.slides.add_slide(slide_layout)

        try:
            # --- Populate Text Placeholders ---

            # Get the name placeholder by its index and set the text
            ph_name = slide.placeholders[PLACEHOLDER_IDX_NAME]
            ph_name.text = str(name)

            # Get the USN placeholder by its index and set the text
            ph_usn = slide.placeholders[PLACEHOLDER_IDX_USN]
            ph_usn.text = str(usn)

            # --- Populate Picture Placeholder ---
            if not os.path.exists(photo_path):
                print(
                    f"    Warning: Photo not found for {name} at '{photo_path}'. Skipping image."
                )
                continue

            # Get the picture placeholder by its index
            ph_photo = slide.placeholders[PLACEHOLDER_IDX_PHOTO]

            # The .insert_picture() method is called on the placeholder itself
            ph_photo.insert_picture(photo_path)

        except KeyError as e:
            print(f"Error: Placeholder index {e} not found on the slide.")
            print("Please double-check your placeholder indices in the config.")
            return
        except Exception as e:
            print(f"Error processing slide for {name}: {e}")

    # Save the final presentation
    try:
        prs.save(OUTPUT_FILE)
        print("\n--- Done! ---")
        print(f"Successfully created '{OUTPUT_FILE}' with {len(df)} student slides.")
    except Exception as e:
        print(f"Error: Could not save the final presentation to '{OUTPUT_FILE}'.")
        print(f"Make sure the file is not open elsewhere. Details: {e}")


if __name__ == "__main__":
    create_presentation()
